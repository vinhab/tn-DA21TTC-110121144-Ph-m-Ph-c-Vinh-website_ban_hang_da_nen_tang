from pptx import Presentation
from pptx.util import Inches, Pt

# Khởi tạo presentation
prs = Presentation()

# Slide 1: Tiêu đề
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Ứng dụng học máy dự đoán nguy cơ mắc bệnh tiểu đường"
slide.placeholders[1].text = "Nhóm thực hiện: ...\nGiảng viên: Thầy Nguyễn Nhứt Lam\nKhoa Kỹ thuật & Công nghệ – Đại học Trà Vinh\nNăm 2025"

# Hàm thêm slide tiêu đề + nội dung
def add_slide(title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    for pt in bullet_points:
        p = body.add_paragraph()
        p.text = pt
        p.level = 0

# Slide 2
add_slide("Lý do chọn đề tài", [
    "Tiểu đường là bệnh lý mãn tính phổ biến, gây biến chứng nguy hiểm.",
    "Số ca mắc tăng nhanh, đặc biệt tại các nước đang phát triển.",
    "Phát hiện sớm giúp điều trị hiệu quả hơn.",
    "Ứng dụng học máy hỗ trợ chẩn đoán và sàng lọc."
])

# Slide 3
add_slide("Mục tiêu nghiên cứu", [
    "Xây dựng mô hình học máy dự đoán tiểu đường từ dữ liệu lâm sàng.",
    "So sánh hiệu quả: Logistic, XGBoost, RF, KNN, SVM.",
    "Tối ưu mô hình và phân tích đặc trưng ảnh hưởng."
])

# Slide 4
add_slide("Đối tượng & Phạm vi", [
    "Dữ liệu: 100.000 cá nhân (Kaggle).",
    "Phân loại nhị phân: mắc / không mắc.",
    "Chỉ sử dụng dữ liệu công khai.",
    "Ứng dụng vào dự đoán y tế dự phòng."
])

# Slide 5
add_slide("Quy trình & Phương pháp", [
    "Tiền xử lý: thiếu, mã hóa, chuẩn hóa.",
    "Trực quan hóa & phân tích dữ liệu.",
    "Xây dựng mô hình & Cross‑validation.",
    "Đánh giá: Accuracy, Precision, Recall, F1, AUC‑ROC."
])

# Slide 6
add_slide("Mô tả bộ dữ liệu", [
    "Nguồn: Comprehensive Diabetes Clinical Dataset (Kaggle, 2023).",
    "15 đặc trưng: Tuổi, BMI, HbA1c, huyết áp, hút thuốc…",
    "Biến mục tiêu: diabetes (1/0)."
])

# Slide 7
add_slide("Tiền xử lý dữ liệu", [
    "Xử lý giá trị thiếu: định lượng → trung vị, phân loại → mode.",
    "One‑hot encoding biến phân loại.",
    "Chuẩn hóa định lượng (StandardScaler).",
    "Cân bằng dữ liệu bằng SMOTE."
])

# Slide 8
add_slide("Các mô hình triển khai", [
    "Logistic Regression",
    "Random Forest",
    "XGBoost",
    "KNN",
    "SVM",
    "Huấn luyện với K‑Fold + GridSearchCV/RandomizedSearchCV."
])

# Slide 9
add_slide("Kết quả hiệu quả mô hình", [
    "XGBoost đạt Accuracy cao nhất: 96.63%.",
    "F1‑score & AUC‑ROC cũng vượt trội.",
    "Các mô hình còn lại thấp hơn."
])

# Slide 10
add_slide("Tối ưu hóa mô hình", [
    "RandomizedSearchCV tối ưu XGBoost.",
    "Accuracy sau tối ưu: 96.73%.",
    "Kiểm thử mẫu mới: hiệu quả tốt, khả năng áp dụng thực tế."
])

# Slide 11
add_slide("Kết luận", [
    "Xây dựng thành công hệ thống dự đoán tiểu đường.",
    "Mô hình học máy cho kết quả khả quan.",
    "XGBoost là mô hình hiệu quả nhất."
])

# Slide 12
add_slide("Hạn chế & Phát triển", [
    "Hạn chế: dữ liệu công khai, chưa đa dạng thực tế.",
    "Hướng phát triển: thêm dữ liệu thực tế.",
    "Thử nghiệm học sâu: ANN, LightGBM, CatBoost.",
    "Phát triển giao diện ứng dụng."
])

# Slide 13
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Tài liệu tham khảo"
body = slide.shapes.placeholders[1].text_frame
body.text = "Priyam Choksi, “100000 Diabetes Clinical Dataset,” Kaggle, 2023."
p = body.add_paragraph()
p.text = "https://www.kaggle.com/datasets/priyamchoksi/100000-diabetes-clinical-dataset"

# Lưu file
prs.save("Demo_Presentation_DuDoan_TieuDuong.pptx")
print("Tạo file PPTX xong: Demo_Presentation_DuDoan_TieuDuong.pptx")
